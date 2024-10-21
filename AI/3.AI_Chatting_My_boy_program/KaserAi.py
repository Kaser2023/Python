import os
import tkinter as tk
from tkinter import filedialog, Text, ttk, messagebox
import google.generativeai as genai
import time
import mimetypes
import docx
import PyPDF2
import openpyxl
import csv
import pptx
import io
import zipfile
import rarfile
from py7zr import SevenZipFile
import PIL.Image
import shutil

# --- YouTube Transcript Extraction ---
from youtube_transcript_api import YouTubeTranscriptApi

# --- Google Gemini API Setup ---
model_names = [
    "gemini-1.5-flash-latest",
    "gemini-1.5-pro-latest",
]


api_key = "AIzaSyAMPldBGV9CeYfSV4MUps6cKtEULNQAydk"

genai.configure(api_key=api_key)

# --- Global Variables ---
chat_sessions = {}  
active_chat_id = None
chat_count = 1  
language_entry = None # 
select_language_button = None 

# --- YouTube Functions ---
def extract_video_id(url):
    """Extracts the video ID from various YouTube URL formats."""
    if "youtu.be/" in url:
        return url.split("youtu.be/")[1].split("?")[0]
    elif "watch?v=" in url:
        return url.split("watch?v=")[1].split("&")[0]
    else:
        raise ValueError("Invalid YouTube URL format.")
    


def fetch_and_start_youtube_chat(video_url):
    """Fetches the YouTube transcript and starts the chat session."""
    global chat_count, chat_text, language_entry, select_language_button

    try:
        # Clear previous language selection UI (if any)
        clear_language_selection()

        video_id = extract_video_id(video_url)

        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        available_languages = [
            (transcript.language_code, transcript.language)
            for transcript in transcript_list
        ]

        # --- Language Selection within the Chat Window ---
        chat_text.config(state="normal") 
        chat_text.insert(tk.END, "Available Transcript Languages:\n")
        for i, (code, lang) in enumerate(available_languages):
            chat_text.insert(tk.END, f"{i + 1}. {lang} ({code})\n")
        chat_text.config(state="disabled")

        def on_language_input():
            global chat_count, language_entry, select_language_button # Access global variables
            try:
                choice = int(language_entry.get()) - 1  
                selected_language_code = available_languages[choice][0]
                
                transcript = transcript_list.find_transcript([selected_language_code])
                transcript_text = transcript.fetch()

                chat_id = f"YouTube Chat {chat_count}"
                chat_count += 1
                create_new_chat_session(
                    chat_id,
                    initial_message=f"I have access to the transcript of a YouTube video. Transcript: {transcript_text}"
                )

                # Clear input area and disable entry
                language_entry.delete(0, tk.END)
                language_entry.config(state="disabled") 
                select_language_button.config(state="disabled")

            except (ValueError, IndexError):
                messagebox.showerror("Error", "Invalid input. Please enter a valid number.")
            except Exception as e:
                messagebox.showerror("Error", f"An error occurred fetching the transcript: {e}")

        # Create the input field and button, initially DISABLED
        language_entry = tk.Entry(chat_frame, state="disabled")  
        language_entry.pack()

        select_language_button = tk.Button(
            chat_frame, text="Select Language (write a number)", command=on_language_input, state="disabled"
        )
        select_language_button.pack()

        # --- Enable language selection UI after successful transcript fetch ---
        language_entry.config(state="normal")
        select_language_button.config(state="normal")

    except ValueError as e:
        messagebox.showerror("Error", str(e))
    except Exception as e:
        if "Subtitles are disabled for this video" in str(e):
            messagebox.showerror("Error", f"The video provided does not have Subtitles!")
            print("The video provided does not have Subtitles")
        else:
            messagebox.showerror("Error", f"An error occurred: {e}")



def clear_language_selection():
    """Removes language selection input and button from the GUI."""
    global language_entry, select_language_button
    if language_entry is not None:
        language_entry.pack_forget() 
        language_entry = None
    if select_language_button is not None:
        select_language_button.pack_forget()
        select_language_button = None 


# --- KaserAI Functions ---

def create_new_chat_session(chat_id, initial_message=None):
    """Creates a new chat session with Gemini and stores it in chat_sessions.
    
    Args:
        chat_id (str): Unique ID for the chat session.
        initial_message (str, optional): An initial message to send to the model.
    """
    global chat_sessions
    for model_name in model_names:
        try:
            generation_config = {
                "temperature": 1, 
                "top_p": 0.95,
            }
            chat_session = genai.GenerativeModel(
                model_name=model_name, generation_config=generation_config
            ).start_chat()
            
            chat_sessions[chat_id] = {
                "session": chat_session,
                "history": [],
                "uploaded_files": []
            }
            
            if initial_message:
                response = chat_session.send_message(initial_message)
                chat_sessions[chat_id]["history"].append(f"Gemini: I have access to the transcript. How can I help you?") 

            update_chat_history_list()
            display_chat(chat_id)
            return True
        except Exception as e:
            print(f"Error trying model '{model_name}': {e}")
    print("Error: Could not create a chat session with any available model.")
    return False

def send_message(event=None):
    """Sends the user's message to Gemini and updates the chat history."""
    global active_chat_id, chat_sessions
    if active_chat_id is None:
        return 

    user_input = input_field.get("1.0", tk.END).strip()
    input_field.delete("1.0", tk.END)

    if user_input:
        chat_data = chat_sessions.get(active_chat_id)
        if chat_data:
            try: 
                chat_session = chat_data["session"]
                prompt_parts = [user_input] + chat_data["uploaded_files"]
                response = chat_session.send_message(prompt_parts)
                
                chat_data["history"].append(f"You: {user_input}")
                chat_data["history"].append(f"Gemini: {response.text}")
                chat_data["uploaded_files"] = []  # Clear uploaded files
                update_chat_display(active_chat_id)

            except Exception as e:
                print(f"Error sending message: {e}")
                messagebox.showerror("Error", f"The Program has finished!")

def start_new_chat():
    """Initiates a new chat session."""
    global chat_count
    chat_id = f"Chat {chat_count}"
    chat_count += 1
    create_new_chat_session(chat_id)



def start_new_youtube_chat():
    """Initiates the process of starting a new YouTube chat within the main window."""
    global chat_text, youtube_url_entry

    def get_url_and_fetch_transcript():
        """Retrieves the URL from the entry field and initiates the fetch process."""
        video_url = youtube_url_entry.get().strip()
        youtube_url_entry.delete(0, tk.END) # Clear the entry field
        youtube_url_entry.pack_forget()  # Hide the entry field
        get_transcript_button.pack_forget()  # Hide the button
        if video_url:  # Proceed only if a URL was entered
            fetch_and_start_youtube_chat(video_url)

    # Create an Entry widget specifically for the YouTube URL
    youtube_url_entry = tk.Entry(chat_frame, width=50)
    youtube_url_entry.pack()

    # Create a button to trigger the transcript fetching
    get_transcript_button = tk.Button(chat_frame, text="Fetch Transcript", command=get_url_and_fetch_transcript)
    get_transcript_button.pack()

    chat_text.config(state="normal")
    chat_text.insert(tk.END, "Enter YouTube Video URL in the field above:\n")
    chat_text.config(state="disabled")


    




# ... (rest of the KaserAI functions remain the same) 


def handle_archive_selection(file_list, chat_data,filename):  # Function is now at a global scope
    """Handles the selection of files from the archive listbox."""
    selected_files = [file_list[i] for i in file_listbox.curselection()]
    if selected_files:
        for file_name in selected_files:
            try:
                temp_dir = os.path.join(".", "temp_extracted")
                os.makedirs(temp_dir, exist_ok=True)
                archive = None  # Ensure archive is defined in this scope
                if filename.endswith(".zip"):
                    archive = zipfile.ZipFile(filename, "r")
                elif filename.endswith(".rar"):
                    archive = rarfile.RarFile(filename, "r")
                elif filename.endswith(".7z"):
                    archive = SevenZipFile(filename, "r")
                if archive:
                    archive.extract(file_name, path=temp_dir)
                    # archive.close()  # Close after extraction
                extracted_file_path = os.path.join(temp_dir, file_name)
                process_file_for_upload(extracted_file_path, chat_data)
            except Exception as e:
                chat_text.config(state="normal")
                chat_text.insert(tk.END, f"Error extracting or analyzing {file_name}: {e}\n")
                chat_text.config(state="disabled")

        shutil.rmtree(temp_dir)
    else:
        messagebox.showwarning("No Files Selected", "Please select at least one file from the archive.")


def on_chat_selected(event):
    """Handles chat selection from the history list."""
    global active_chat_id, select_button  # Make select_button globally accessible

    selection = chat_history_list.curselection()
    if selection:
        index = selection[0]
        active_chat_id = chat_history_list.get(index)
        display_chat(active_chat_id)

        file_listbox.delete(0, tk.END)

        # Try to get the existing select_button
        select_button = file_list_frame.winfo_children()
        select_button = select_button[0] if select_button else None

        chat_data = chat_sessions.get(active_chat_id)
        if chat_data and "archive_file_list" in chat_data:
            
            for file_name in chat_data["archive_file_list"]:
                file_listbox.insert(tk.END, file_name)

            if "archive_filename" in chat_data:
                filename = chat_data["archive_filename"]
                if select_button:
                    # Reuse the existing button by updating its command
                    select_button.config(
                        command=lambda f=filename, c=chat_data: handle_archive_selection(
                            c["archive_file_list"], c, f
                        )
                    )
                else:
                    # Create a new button if it doesn't exist
                    select_button = tk.Button(
                        file_list_frame,
                        text="Select from Archive",
                        command=lambda f=filename, c=chat_data: handle_archive_selection(
                            c["archive_file_list"], c, f
                        ),
                    )
                    select_button.pack()



def display_chat(chat_id):
    """Updates the chat area with the selected chat's history."""
    global chat_sessions
    chat_text.config(state="normal")
    chat_text.delete("1.0", tk.END)

    chat_data = chat_sessions.get(chat_id)
    if chat_data:
        for message in chat_data["history"]:
            chat_text.insert(tk.END, f"{message}\n")
    chat_text.config(state="disabled")

def update_chat_display(chat_id):
    """Refreshes the chat display after sending a message or uploading a file."""
    global chat_sessions 
    chat_text.config(state="normal")
    chat_text.delete("1.0", tk.END) 
    chat_data = chat_sessions.get(chat_id)
    if chat_data:
        for message in chat_data["history"]:
            chat_text.insert(tk.END, f"{message}\n")
    chat_text.config(state="disabled") 

def update_chat_history_list():
    """Refreshes the list of active chats in the GUI."""
    chat_history_list.delete(0, tk.END)
    for i, chat_id in enumerate(chat_sessions.keys()):
        chat_history_list.insert(tk.END, chat_id)
        if i == 0:  # Select the first chat by default
            chat_history_list.selection_set(0)
            on_chat_selected(None)  # Load the first chat initially


def analyze_files(filename):
    global uploaded_files
    try:
        file_extension = os.path.splitext(filename)[1].lower()
        if file_extension == ".docx":
            doc = docx.Document(filename)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif file_extension == ".pdf":
            with open(filename, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                return "".join(
                    pdf_reader.pages[page].extract_text()
                    for page in range(len(pdf_reader.pages))
                )
        elif file_extension == ".csv":
            with open(filename, "r", encoding="utf-8") as file:
                reader = csv.reader(file)
                return "\n".join(", ".join(row) for row in reader)
        elif file_extension in [".pptx", ".ppt"]:
            prs = pptx.Presentation(filename)
            images = [
                PIL.Image.open(io.BytesIO(shape.image.blob))
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE
            ]
            return "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
        elif file_extension == ".xlsx":
            workbook = openpyxl.load_workbook(filename)
            worksheet = workbook.active
            return "\n".join(
                str(cell.value) for row in worksheet.iter_rows() for cell in row
            )
        elif file_extension == ".js":
            with open(filename, "r", encoding="utf-8") as file:
                return file.read()
            
        elif file_extension == ".java":
            with open(filename, "r", encoding="utf-8") as file:
                return file.read()
    except (FileNotFoundError, Exception) as e:
        print(f"Error: {e}. Please enter a valid file path.")
        return ""

def upload_files():
    """Handles file uploads and associates them with the active chat."""
    global active_chat_id, chat_sessions
    if active_chat_id is None:
        return

    filenames = filedialog.askopenfilenames(
        initialdir="/", title="Select files", filetypes=(("All files", "*.*"),)
    )
    if filenames:
        chat_data = chat_sessions[active_chat_id]
        for filename in filenames:
            process_file_for_upload(filename, chat_data)

def process_file_for_upload(filename, chat_data):
    """Processes a single file, handling archives and individual files."""
    try:
        file_extension = os.path.splitext(filename)[1].lower()

        if file_extension in (".zip", ".rar", ".7z"):
            analyze_compressed_file(filename, chat_data)
        elif file_extension in (".pdf", ".docx", ".csv", ".pptx", ".ppt", ".xlsx", ".js", ".java"):
            content = analyze_files(filename)
            chat_data["uploaded_files"].append(content)
            chat_text.config(state="normal")
            chat_text.insert(tk.END, f" Uploaded file is Ready!\n")
            chat_text.config(state="disabled")
        else:
            upload_single_file_to_gemini(filename, chat_data)

    except Exception as e:
        chat_text.config(state="normal")
        chat_text.insert(tk.END, f"Error uploading file {filename}: {e}\n")
        chat_text.config(state="disabled")


def analyze_compressed_file(filename, chat_data):
    """Analyzes files within a compressed archive (zip, rar, 7z)."""
    global active_chat_id, chat_sessions

    try:
        if filename.endswith(".zip"):
            archive = zipfile.ZipFile(filename, "r")
        elif filename.endswith(".rar"):
            archive = rarfile.RarFile(filename, "r")
        elif filename.endswith(".7z"):
            archive = SevenZipFile(filename, "r")
        else:
            messagebox.showerror("Error", "Unsupported archive format. Please choose a zip, rar, or 7z file.")
            return


        file_list = archive.namelist()
        chat_data["archive_file_list"] = file_list

        # Clear the file list and remove any existing button
        file_listbox.delete(0, tk.END)
        for widget in file_list_frame.winfo_children():
            if isinstance(widget, tk.Button):
                widget.destroy()

        for file_name in file_list:
            file_listbox.insert(tk.END, file_name)

        # Create the button ONLY if it doesn't exist
        button_exists = any(isinstance(widget, tk.Button) for widget in file_list_frame.winfo_children())
        if not button_exists:
            select_button = tk.Button(
                file_list_frame,
                text="Select from Archive",
                command=lambda f=filename, c=chat_data: handle_archive_selection(
                    c["archive_file_list"], c, f
                ),
            )
            select_button.pack()  # Pack the button here
        
    except Exception as e:
        chat_text.config(state="normal")
        chat_text.insert(tk.END, f"Error processing archive: {e}\n")
        chat_text.config(state="disabled")


def upload_single_file_to_gemini(filename, chat_data):
    """Uploads a single file to Gemini and updates the chat."""
    try:
        uploaded_file = genai.upload_file(filename)
        processed_file = wait_for_file_processing(uploaded_file)
        chat_data["uploaded_files"].append(processed_file)
        chat_text.config(state="normal")
        chat_text.insert(tk.END, f" Uploaded file is Ready!\n")
        chat_text.config(state="disabled")
        print(f"Uploaded file '{uploaded_file.display_name}' as: {uploaded_file.uri}")
        print(uploaded_file.mime_type)
        print(uploaded_file.name)

    except Exception as e:
        chat_text.config(state="normal")
        chat_text.insert(tk.END, f"Error uploading file {filename}: {e}\n")
        chat_text.config(state="disabled")

def wait_for_file_processing(file):
    """Waits for a file to finish processing on Gemini's servers."""
    while file.state.name == "PROCESSING":
        time.sleep(2)
        file = genai.get_file(file.name)
    return file


def clear_file_list():
    """Clears the file listbox and destroys any select button."""
    file_listbox.delete(0, tk.END)
    for widget in file_list_frame.winfo_children():
        if isinstance(widget, tk.Button):
            widget.destroy()



# --- GUI Setup ---
window = tk.Tk()
window.title("Kaser-Ai")

# --- Left Side: Chat History ---
history_frame = tk.Frame(window, width=200, borderwidth=1, relief="solid")
history_frame.pack(side="left", fill="y")

new_chat_button = tk.Button(history_frame, text="New Chat", command=start_new_chat)
new_chat_button.pack(pady=5)

new_youtube_chat_button = tk.Button(history_frame, text="New YouTube Chat", command=start_new_youtube_chat)
new_youtube_chat_button.pack(pady=5)

chat_history_list = tk.Listbox(history_frame)
chat_history_list.pack(fill="both", expand=True)
chat_history_list.bind('<<ListboxSelect>>', on_chat_selected)


# --- Right Side: Chat Area ---
right_frame = tk.Frame(window)
right_frame.pack(side="left", fill="both", expand=True)

# Chat Area (in the right frame)
chat_frame = tk.Frame(right_frame)
chat_frame.pack(side="top", fill="both", expand=True)


# chat_frame = tk.Frame(window)
# chat_frame.pack(side="left", fill="both", expand=True)

# Chat Output
chat_text = Text(chat_frame, height=20, width=80)
chat_text.config(state="disabled")
chat_text.pack(fill="both", expand=True)

# File Upload
upload_button = tk.Button(chat_frame, text="Upload Files", command=upload_files)
upload_button.pack()

# User Input
input_field = Text(chat_frame, height=3, width=80)
input_field.pack(fill="x")
input_field.bind("<Return>", send_message)  # Send on Enter key

# Send Button
send_button = tk.Button(chat_frame, text="Send", command=send_message)
send_button.pack()

# File List Area (in the right frame, below chat)
file_list_frame = tk.Frame(right_frame, borderwidth=1, relief="solid")
file_list_frame.pack(side="bottom", fill="x")

file_listbox = tk.Listbox(file_list_frame, selectmode="multiple")
file_listbox.pack(fill="both", expand=True)


# ----- Start with new chat -----
start_new_chat()

window.mainloop()